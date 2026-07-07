import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 2. Design System Tokens (Color Palette)
BG_COLOR = RGBColor(15, 23, 42)      # Slate 900 (Deep navy-slate)
CARD_BG = RGBColor(30, 41, 59)        # Slate 800 (Card background)
CARD_BORDER = RGBColor(51, 65, 85)    # Slate 700 (Card border)
ACCENT_BLUE = RGBColor(14, 165, 233)  # Sky 500 (Primary accent)
ACCENT_GREEN = RGBColor(16, 185, 129) # Emerald 500 (Success/Secure accent)
ACCENT_RED = RGBColor(244, 63, 94)    # Rose 500 (Risk/Tamper accent)
TEXT_WHITE = RGBColor(248, 250, 252)  # Slate 50 (Primary text)
TEXT_MUTED = RGBColor(148, 163, 184)  # Slate 400 (Secondary text)

FONT_HEADING = "Trebuchet MS"
FONT_BODY = "Calibri"

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def draw_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_title(slide, text):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_HEADING
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    # Add a thin accent line under the title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(2.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

def add_bullet_points(tf, points, font_size=13, color=TEXT_WHITE):
    for i, pt in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        p.space_after = Pt(6)
        
        if isinstance(pt, tuple):
            text, is_bold, indent, *rest = pt
            custom_color = rest[0] if rest else color
        else:
            text = pt
            is_bold = False
            indent = 0
            custom_color = color
            
        p.level = indent
        p.text = ("• " if indent == 0 else "  - ") + text
        p.font.name = FONT_BODY
        p.font.size = Pt(font_size - (indent * 1.5))
        p.font.bold = is_bold
        p.font.color.rgb = custom_color

slide_layout = prs.slide_layouts[6] # blank layout

# ----------------- SLIDE 1: COVER -----------------
slide1 = prs.slides.add_slide(slide_layout)
set_slide_background(slide1)
draw_card(slide1, 0.8, 1.0, 11.73, 5.5, bg_color=CARD_BG, border_color=CARD_BORDER)

# Title Text inside card
title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(10.93), Inches(1.8))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "RANCANG BANGUN SISTEM PENGAMANAN PESAN BERBASIS\nHYBRID CRYPTOGRAPHY AES-RSA MENGGUNAKAN\nALGORITMA NAIVE BAYES"
p.font.name = FONT_HEADING
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Aplikasi Chatting E2EE 'ChatKu' dengan Deteksi Risiko Lokal (Client-Side)"
p2.font.name = FONT_BODY
p2.font.size = Pt(15)
p2.font.color.rgb = ACCENT_BLUE
p2.space_before = Pt(8)
p2.alignment = PP_ALIGN.CENTER

# Metadata (Author)
meta_box = slide1.shapes.add_textbox(Inches(1.2), Inches(3.2), Inches(5.3), Inches(2.2))
tf_meta = meta_box.text_frame
tf_meta.word_wrap = True
add_bullet_points(tf_meta, [
    ("Oleh Peneliti:", True, 0, ACCENT_BLUE),
    ("Dimas Kurniawan", True, 0, TEXT_WHITE),
    ("NIM: 2022573010009", False, 0, TEXT_MUTED),
    ("Program Studi: Teknik Informatika", False, 0, TEXT_MUTED),
    ("Politeknik Negeri Lhokseumawe", False, 0, TEXT_MUTED)
], font_size=12)

# Metadata (Advisors & Examiners)
advisor_box = slide1.shapes.add_textbox(Inches(6.8), Inches(3.2), Inches(5.3), Inches(3.0))
tf_adv = advisor_box.text_frame
tf_adv.word_wrap = True
add_bullet_points(tf_adv, [
    ("Komisi Sidang & Pembimbing:", True, 0, ACCENT_BLUE),
    ("Pembimbing I / Ketua Sidang:", True, 0, TEXT_WHITE),
    ("M. Khadafi, S.T., M.T.", False, 1, TEXT_WHITE),
    ("Pembimbing II / Sekretaris Sidang:", True, 0, TEXT_WHITE),
    ("Muhammad Reza Zulman, S.S.T., M.Sc.", False, 1, TEXT_WHITE),
    ("Dosen Penguji:", True, 0, ACCENT_BLUE),
    ("Penguji I: Salahuddin, S.T., M.Cs.", False, 1, TEXT_WHITE),
    ("Penguji II: Zulfan Khairil Simbolon, S.T., M.Eng", False, 1, TEXT_WHITE),
    ("Penguji III: Mahdi, S.T., M.Cs.", False, 1, TEXT_WHITE)
], font_size=11)


# ----------------- SLIDE 2: LATAR BELAKANG -----------------
slide2 = prs.slides.add_slide(slide_layout)
set_slide_background(slide2)
add_title(slide2, "Latar Belakang & Urgensi Penelitian")

# Column 1 Card (Problem)
draw_card(slide2, 0.8, 1.6, 5.6, 5.1, bg_color=CARD_BG, border_color=ACCENT_RED)
problem_title = slide2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(0.4))
problem_title.text_frame.word_wrap = True
p = problem_title.text_frame.paragraphs[0]
p.text = "PERMASALAHAN UTAMA"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_RED

problem_box = slide2.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.2))
tf_prob = problem_box.text_frame
tf_prob.word_wrap = True
add_bullet_points(tf_prob, [
    ("Tingginya volume komunikasi data sensitif (personal, finansial, kredensial) pada aplikasi pesan instan.", False, 0),
    ("Maraknya ancaman siber berupa phishing, penyebaran malware via file APK palsu, dan social engineering di Indonesia.", False, 0),
    ("Enkripsi E2EE standar mengamankan data secara agnostik, seluruh pesan baik yang biasa maupun yang berisi ancaman/penipuan diamankan secara sama sehingga penerima tetap menerima konten berbahaya.", False, 0),
    ("Proses klasifikasi risiko di sisi server (Server-Side Decryption) melanggar privasi mutlak pengguna karena server harus membuka plaintext pesan asli.", False, 0)
], font_size=13)

# Column 2 Card (Solution)
draw_card(slide2, 6.9, 1.6, 5.6, 5.1, bg_color=CARD_BG, border_color=ACCENT_GREEN)
sol_title = slide2.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(0.4))
sol_title.text_frame.word_wrap = True
p = sol_title.text_frame.paragraphs[0]
p.text = "SOLUSI YANG DIUSULKAN"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN

sol_box = slide2.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2))
tf_sol = sol_box.text_frame
tf_sol.word_wrap = True
add_bullet_points(tf_sol, [
    ("Client-Side Classification: Deteksi tingkat risiko pesan otomatis dijalankan lokal di RAM HP pengirim menggunakan model Naive Bayes.", False, 0),
    ("Fallback Content Replacement: Pesan terdeteksi Berisiko diganti menjadi pesan fallback lokal ('Pesan ini terindikasi berisiko') sebelum enkripsi.", False, 0),
    ("Hybrid Cryptography AES-RSA: Pesan aman dienkripsi end-to-end dengan AES-256-CBC, sedangkan kunci AES dilindungi dengan RSA-2048.", False, 0),
    ("Zero-Knowledge & Zero-Storage Server: Server FastAPI dan MongoDB murni bertindak sebagai router buta data yang hanya menyimpan ciphertext acak.", False, 0)
], font_size=13)


# ----------------- SLIDE 3: RUMUSAN & BATASAN -----------------
slide3 = prs.slides.add_slide(slide_layout)
set_slide_background(slide3)
add_title(slide3, "Rumusan Masalah & Batasan Penelitian")

# Rumusan Masalah
draw_card(slide3, 0.8, 1.6, 5.6, 5.1, bg_color=CARD_BG)
rumusan_title = slide3.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(0.4))
p = rumusan_title.text_frame.paragraphs[0]
p.text = "RUMUSAN MASALAH"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

rumusan_box = slide3.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.2))
tf_rum = rumusan_box.text_frame
tf_rum.word_wrap = True
add_bullet_points(tf_rum, [
    ("1. Bagaimana merancang sistem pengamanan pesan berbasis Hybrid Cryptography AES-RSA yang dikombinasikan dengan Naive Bayes untuk klasifikasi risiko?", False, 0),
    ("2. Bagaimana tingkat akurasi algoritma Naive Bayes dalam mengklasifikasikan tingkat risiko data pesan secara offline?", False, 0),
    ("3. Bagaimana hasil pengujian keamanan dan fungsionalitas sistem yang dirancang?", False, 0)
], font_size=13)

# Batasan Masalah
draw_card(slide3, 6.9, 1.6, 5.6, 5.1, bg_color=CARD_BG)
batasan_title = slide3.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(0.4))
p = batasan_title.text_frame.paragraphs[0]
p.text = "BATASAN MASALAH"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

batasan_box = slide3.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2))
tf_bat = batasan_box.text_frame
tf_bat.word_wrap = True
add_bullet_points(tf_bat, [
    ("1. Pembahasan berfokus pada perancangan & implementasi pengamanan data pesan, bukan keamanan hardware atau jaringan secara menyeluruh.", False, 0),
    ("2. Kriptografi hybrid dibatasi pada penggabungan algoritma AES (enkripsi pesan) dan RSA (enkripsi kunci).", False, 0),
    ("3. Proses enkripsi dan pengamanan hanya diterapkan pada pesan dalam bentuk teks.", False, 0),
    ("4. Algoritma Naive Bayes digunakan hanya untuk klasifikasi pesan berdasarkan tingkat keamanannya.", False, 0),
    ("5. Klasifikasi tingkat keamanan data dibatasi pada kategori Berisiko dan Tidak Berisiko.", False, 0),
    ("6. Pengujian keamanan dibatasi pada simulasi serangan pengiriman data, bukan penetration testing menyeluruh.", False, 0)
], font_size=13)


# ----------------- SLIDE 4: ARSITEKTUR SISTEM -----------------
slide4 = prs.slides.add_slide(slide_layout)
set_slide_background(slide4)
add_title(slide4, "Arsitektur & Alur Kerja Sistem ChatKu")

# 3 Horizontal Cards represent layers
# Card 1: Input & Preprocess
draw_card(slide4, 0.8, 1.8, 3.6, 4.8, bg_color=CARD_BG, border_color=CARD_BORDER)
c1_title = slide4.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.2), Inches(0.4))
p = c1_title.text_frame.paragraphs[0]
p.text = "1. PRAPEMROSESAN LOKAL"
p.font.name = FONT_HEADING
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

c1_box = slide4.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(3.2), Inches(3.9))
tf_c1 = c1_box.text_frame
tf_c1.word_wrap = True
add_bullet_points(tf_c1, [
    ("Pengirim mengetik pesan teks di UI aplikasi.", False, 0),
    ("Case Folding: Mengubah huruf menjadi lowercase.", False, 0),
    ("Pattern Replacement: Regex mengganti pola sensitif (URL, nomor HP, OTP, nominal uang) menjadi token khusus.", False, 0),
    ("Tokenisasi & Stopword removal.", False, 0),
    ("Stemming Sastrawi: Mengembalikan kata berimbuhan ke kata dasar Bahasa Indonesia.", False, 0)
], font_size=12)

# Card 2: AI Classification
draw_card(slide4, 4.8, 1.8, 3.6, 4.8, bg_color=CARD_BG, border_color=ACCENT_BLUE)
c2_title = slide4.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(3.2), Inches(0.4))
p = c2_title.text_frame.paragraphs[0]
p.text = "2. DETEKSI RISIKO (CLIENT)"
p.font.name = FONT_HEADING
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

c2_box = slide4.shapes.add_textbox(Inches(5.0), Inches(2.5), Inches(3.2), Inches(3.9))
tf_c2 = c2_box.text_frame
tf_c2.word_wrap = True
add_bullet_points(tf_c2, [
    ("Model Naive Bayes + Safety Net (Daftar Kata Kasar/Obfuscation) di perangkat mengevaluasi teks bersih.", False, 0),
    ("Keputusan Pengiriman:", True, 0, ACCENT_BLUE),
    ("Jika Berisiko: Teks asli diganti lokal dengan pesan fallback: 'Pesan ini terindikasi berisiko'.", False, 1, ACCENT_RED),
    ("Jika Aman: Pesan asli dipertahankan.", False, 1, ACCENT_GREEN),
    ("Pembersihan RAM: Data sensitif dihapus dari RAM setelah klasifikasi selesai.", False, 0)
], font_size=12)

# Card 3: Cryptography & Transmission
draw_card(slide4, 8.8, 1.8, 3.6, 4.8, bg_color=CARD_BG, border_color=CARD_BORDER)
c3_title = slide4.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.2), Inches(0.4))
p = c3_title.text_frame.paragraphs[0]
p.text = "3. HYBRID ENCRYPTION"
p.font.name = FONT_HEADING
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

c3_box = slide4.shapes.add_textbox(Inches(9.0), Inches(2.5), Inches(3.2), Inches(3.9))
tf_c3 = c3_box.text_frame
tf_c3.word_wrap = True
add_bullet_points(tf_c3, [
    ("Kunci AES-256 dibangkitkan acak per pesan.", False, 0),
    ("Pesan dienkripsi AES-256-CBC.", False, 0),
    ("Kunci AES dienkripsi dengan RSA-2048 menggunakan kunci publik Penerima (ek_user) dan kunci publik Pengirim (ek_sender).", False, 0),
    ("Plaintext di-hash SHA-256.", False, 0),
    ("Payload terenkripsi dikirim via WebSocket ke MongoDB. Server murni menyimpan ciphertext.", False, 0)
], font_size=12)


# ----------------- SLIDE 5: DATASET -----------------
slide5 = prs.slides.add_slide(slide_layout)
set_slide_background(slide5)
add_title(slide5, "Data, Sumber Dataset & Preprocessing")

# Left Column (Sources & Split)
draw_card(slide5, 0.8, 1.6, 5.6, 5.1, bg_color=CARD_BG)
ds_title = slide5.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(0.4))
p = ds_title.text_frame.paragraphs[0]
p.text = "SUMBER DATA & PEMBAGIAN DATA"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

ds_box = slide5.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.2))
tf_ds = ds_box.text_frame
tf_ds.word_wrap = True
add_bullet_points(tf_ds, [
    ("Dataset Awal: 38.070 sampel pesan Bahasa Indonesia.", True, 0, ACCENT_BLUE),
    ("Pembersihan data (Data Cleaning) membuang 1.700 sampel tidak layak. Sisa: 37.130 sampel final.", False, 1),
    ("Distribusi Kelas (Tabel 4.1):", True, 0, ACCENT_BLUE),
    ("Berisiko: 4.196 sampel (11,0%)", False, 1),
    ("Tidak Berisiko: 33.890 sampel (89,0%)", False, 1),
    ("Rasio ketidakseimbangan kelas (imbalance ratio) sebesar 1:10,2.", False, 1),
    ("Pembagian Subset Data (70:15:15):", True, 0, ACCENT_BLUE),
    ("Data Latih: 25.991 sampel (70%)", False, 1),
    ("Data Validasi: 5.394 sampel (15%)", False, 1),
    ("Data Uji: 5.394 sampel (15%)", False, 1)
], font_size=12)

# Right Column (Preprocessing Detail)
draw_card(slide5, 6.9, 1.6, 5.6, 5.1, bg_color=CARD_BG)
prep_title = slide5.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(0.4))
p = prep_title.text_frame.paragraphs[0]
p.text = "TAHAHPAN PREPROCESSING LOKAL"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

prep_box = slide5.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2))
tf_prep = prep_box.text_frame
tf_prep.word_wrap = True
add_bullet_points(tf_prep, [
    ("Case Folding: Mengubah seluruh karakter menjadi huruf kecil.", False, 0),
    ("Pattern Replacement (Regex): Mengubah pola khusus:", True, 0, ACCENT_BLUE),
    ("URL/Link mencurigakan → token 'URL_CURIGA'", False, 1),
    ("Nomor Telepon/HP asing → token 'NOMOR_HP_ASING'", False, 1),
    ("Kode OTP (5-8 digit) → token 'KODE_OTP'", False, 1),
    ("Nominal Uang (Rp...) → token 'NOMINAL_UANG'", False, 1),
    ("Tokenisasi: Memecah kalimat menjadi token kata berdasarkan spasi.", False, 0),
    ("Stopword Removal: Menghapus kata umum Bahasa Indonesia (yang, di, ke, dll.) dengan mengecualikan kata kasar.", False, 0),
    ("Stemming Sastrawi: Mengubah kata berimbuhan menjadi kata dasar Bahasa Indonesia.", False, 0)
], font_size=12)


# ----------------- SLIDE 6: KRIPTOGRAFI HYBRID -----------------
slide6 = prs.slides.add_slide(slide_layout)
set_slide_background(slide6)
add_title(slide6, "Skema Kriptografi Hybrid & Integritas Data")

# Card 1: AES
draw_card(slide6, 0.8, 1.8, 3.6, 4.8, bg_color=CARD_BG, border_color=CARD_BORDER)
cr1_title = slide6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.2), Inches(0.4))
p = cr1_title.text_frame.paragraphs[0]
p.text = "ENKRIPSI DATA: AES-256-CBC"
p.font.name = FONT_HEADING
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

cr1_box = slide6.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(3.2), Inches(3.9))
tf_cr1 = cr1_box.text_frame
tf_cr1.word_wrap = True
add_bullet_points(tf_cr1, [
    ("Menggunakan Advanced Encryption Standard (AES) simetris dengan kunci 256-bit.", False, 0),
    ("Mode operasi CBC (Cipher Block Chaining) dengan Initialization Vector (IV) 16-byte acak.", False, 0),
    ("Penerapan One-Time Key: Setiap pesan baru membangkitkan kunci AES & IV acak, sehingga ciphertext selalu berbeda meski isi pesan sama.", False, 0),
    ("Padding PKCS#7 menyelaraskan panjang pesan kelipatan 16-byte.", False, 0)
], font_size=11)

# Card 2: RSA Dual-Key
draw_card(slide6, 4.8, 1.8, 3.6, 4.8, bg_color=CARD_BG, border_color=ACCENT_GREEN)
cr2_title = slide6.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(3.2), Inches(0.4))
p = cr2_title.text_frame.paragraphs[0]
p.text = "DISTRIBUSI KUNCI: RSA-2048"
p.font.name = FONT_HEADING
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN

cr2_box = slide6.shapes.add_textbox(Inches(5.0), Inches(2.5), Inches(3.2), Inches(3.9))
tf_cr2 = cr2_box.text_frame
tf_cr2.word_wrap = True
add_bullet_points(tf_cr2, [
    ("Kunci AES dilindungi secara asimetris dengan RSA-2048 (padding PKCS#1 v1.5 / OAEP).", False, 0),
    ("Skema Dual-Key Encryption:", True, 0, ACCENT_BLUE),
    ("ek_user: Kunci AES dienkripsi dengan kunci publik Penerima agar hanya penerima yang bisa membaca.", False, 1),
    ("ek_sender: Kunci AES dienkripsi dengan kunci publik Pengirim agar pengirim bisa membaca riwayat chat sendiri.", False, 1),
    ("Kunci privat disimpan di perangkat klien, server tidak pernah menerima kunci privat.", False, 0)
], font_size=11)

# Card 3: Integrity
draw_card(slide6, 8.8, 1.8, 3.6, 4.8, bg_color=CARD_BG, border_color=CARD_BORDER)
cr3_title = slide6.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.2), Inches(0.4))
p = cr3_title.text_frame.paragraphs[0]
p.text = "SEGAL HASH & DB STORAGE"
p.font.name = FONT_HEADING
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

cr3_box = slide6.shapes.add_textbox(Inches(9.0), Inches(2.5), Inches(3.2), Inches(3.9))
tf_cr3 = cr3_box.text_frame
tf_cr3.word_wrap = True
add_bullet_points(tf_cr3, [
    ("SHA-256 Integrity Verification: Nilai hash SHA-256 dari plaintext asli dihitung sebelum enkripsi.", False, 0),
    ("Hash dikirim bersama ciphertext dan diverifikasi ulang pasca dekripsi di penerima.", False, 0),
    ("Struktur database MongoDB hanya menyimpan: ciphertext, IV, ek_user, ek_sender, dan hash SHA-256. Server murni buta data (Zero-Knowledge).", False, 0)
], font_size=11)


# ----------------- SLIDE 7: IMPLEMENTASI & UI -----------------
slide7 = prs.slides.add_slide(slide_layout)
set_slide_background(slide7)
add_title(slide7, "Implementasi Sistem & Antarmuka Aplikasi ChatKu")

# Left Column (Tech Stack)
draw_card(slide7, 0.8, 1.6, 5.6, 5.1, bg_color=CARD_BG)
tech_title = slide7.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(0.4))
p = tech_title.text_frame.paragraphs[0]
p.text = "STACK TEKNOLOGI SISTEM"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

tech_box = slide7.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.2))
tf_tech = tech_box.text_frame
tf_tech.word_wrap = True
add_bullet_points(tf_tech, [
    ("Frontend Mobile: React Native (TypeScript)", True, 0, ACCENT_BLUE),
    ("Library Kriptografi: node-forge (RSA-OAEP), react-native-aes-crypto (AES-256-CBC).", False, 1),
    ("Stemming Bahasa Indonesia: Library Sastrawi JS yang berjalan lokal di RAM HP pengirim.", False, 1),
    ("Backend FastAPI: Python-based web server berkinerja tinggi (ASGI).", True, 0, ACCENT_BLUE),
    ("Protokol WebSocket real-time untuk transmisi payload terenkripsi secara instan.", False, 1),
    ("Database MongoDB: Penyimpanan NoSQL berkinerja tinggi untuk user, pesan terenkripsi, dan forensic log.", True, 0, ACCENT_BLUE)
], font_size=13)

# Right Column (UI Screens)
draw_card(slide7, 6.9, 1.6, 5.6, 5.1, bg_color=CARD_BG)
ui_title = slide7.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(0.4))
p = ui_title.text_frame.paragraphs[0]
p.text = "FITUR ANTARMUKA UTAMA (UI)"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

ui_box = slide7.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2))
tf_ui = ui_box.text_frame
tf_ui.word_wrap = True
add_bullet_points(tf_ui, [
    ("Onboarding Page: Memandu pengguna memahami E2EE, Zero-Storage, dan Anti-Forensik.", False, 0),
    ("Chatroom Page: Menampilkan chat secara real-time. Jika sensor AI mendeteksi pesan berisiko kasar, pengiriman langsung diblokir secara lokal dan diganti pesan fallback.", False, 0),
    ("Security Center Dashboard:", True, 0, ACCENT_BLUE),
    ("Dashboard: Menampilkan skor keamanan dan fitur aktif.", False, 1),
    ("Key Verify: Menampilkan fingerprint kunci publik RSA untuk memverifikasi kontak guna mencegah serangan MITM.", False, 1),
    ("Forensic Log: Riwayat aktivitas keamanan (auth, crypto, access) yang dilindungi SHA-256.", False, 1)
], font_size=12)


# ----------------- SLIDE 8: ACCURACY RESULTS -----------------
slide8 = prs.slides.add_slide(slide_layout)
set_slide_background(slide8)
add_title(slide8, "Hasil Pengujian & Akurasi Model Naive Bayes")

# Left Column (Metrics)
draw_card(slide8, 0.8, 1.6, 5.6, 5.1, bg_color=CARD_BG, border_color=ACCENT_BLUE)
acc_title = slide8.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(0.4))
p = acc_title.text_frame.paragraphs[0]
p.text = "METRIK EVALUASI MODEL"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

acc_box = slide8.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.2))
tf_acc = acc_box.text_frame
tf_acc.word_wrap = True
add_bullet_points(tf_acc, [
    ("Model Klasifikasi: Multinomial Naive Bayes", True, 0, ACCENT_BLUE),
    ("Alpha = 1.0 (Laplace Smoothing) | Fitur: TF-IDF N-Gram (1,2) dengan 9.847 kata unik.", False, 1),
    ("Akurasi Rata-rata (Weighted Avg): 92,90%", True, 0, ACCENT_GREEN),
    ("Evaluasi Per Kelas pada Data Uji (Tabel 4.6):", True, 0, ACCENT_BLUE),
    ("Kelas Berisiko (Support = 581):", True, 1, TEXT_WHITE),
    ("Precision: 66,45%  |  Recall: 68,85% (91,5% pada Bab IV/V text)  |  F1-Score: 67,62%", False, 2, TEXT_MUTED),
    ("Kelas Tidak Berisiko (Support = 4.813):", True, 1, TEXT_WHITE),
    ("Precision: 95,80%  |  Recall: 95,80%  |  F1-Score: 95,80%", False, 2, TEXT_MUTED),
    ("Catatan: Nilai recall/precision seimbang diperkuat dengan threshold klasifikasi sebesar 0,65.", False, 0)
], font_size=11)

# Right Column (Confusion Matrix)
draw_card(slide8, 6.9, 1.6, 5.6, 5.1, bg_color=CARD_BG)
cm_title = slide8.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(0.4))
p = cm_title.text_frame.paragraphs[0]
p.text = "CONFUSION MATRIX (DATA UJI)"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

cm_box = slide8.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2))
tf_cm = cm_box.text_frame
tf_cm.word_wrap = True
add_bullet_points(tf_cm, [
    ("Matriks Evaluasi Data Uji (5.394 total sampel):", True, 0, ACCENT_BLUE),
    ("TP (True Positive) = 400 sampel", True, 0, ACCENT_GREEN),
    ("Pesan berisiko berhasil terdeteksi dan diblokir lokal.", False, 1, TEXT_MUTED),
    ("TN (True Negative) = 4.611 sampel", True, 0, ACCENT_GREEN),
    ("Pesan aman diizinkan lolos untuk proses enkripsi.", False, 1, TEXT_MUTED),
    ("FP (False Positive) = 202 sampel", True, 0, ACCENT_RED),
    ("Pesan aman salah terdeteksi berisiko (false alarm).", False, 1, TEXT_MUTED),
    ("FN (False Negative) = 181 sampel", True, 0, ACCENT_RED),
    ("Pesan berisiko lolos dari sensor model Naive Bayes.", False, 1, TEXT_MUTED),
    ("Mitigasi FN: Safety Net (rule-based filtering) bertindak sebagai penyaring lapis pertama untuk memblokir kata kasar/obfuscation sensitif secara mutlak.", True, 0, ACCENT_BLUE)
], font_size=11)


# ----------------- SLIDE 9: SECURITY TESTING -----------------
slide9 = prs.slides.add_slide(slide_layout)
set_slide_background(slide9)
add_title(slide9, "Hasil Pengujian Keamanan & Integritas Sistem")

# Left Column (Network & DB Audit)
draw_card(slide9, 0.8, 1.6, 5.6, 5.1, bg_color=CARD_BG)
net_title = slide9.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(0.4))
p = net_title.text_frame.paragraphs[0]
p.text = "AUDIT JARINGAN & BASIS DATA"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

net_box = slide9.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.2))
tf_net = net_box.text_frame
tf_net.word_wrap = True
add_bullet_points(tf_net, [
    ("Penyadapan Jaringan (Wireshark MITM):", True, 0, ACCENT_BLUE),
    ("Paket WebSocket yang disadap di tengah jalan hanya menampilkan payload Base64 terenkripsi (ciphertext).", False, 1),
    ("Tidak ada plaintext pesan atau kunci AES mentah yang bocor di jaringan transit.", False, 1),
    ("Zero-Knowledge Database Audit:", True, 0, ACCENT_BLUE),
    ("Isi database MongoDB diperiksa dan terbukti hanya menyimpan payload terenkripsi (ciphertext), IV, ek_user, ek_sender, dan nilai hash.", False, 1),
    ("Server dan database tidak menyimpan teks pesan asli maupun kunci privat pengguna.", False, 1)
], font_size=13)

# Right Column (Integrity & Testing)
draw_card(slide9, 6.9, 1.6, 5.6, 5.1, bg_color=CARD_BG)
int_title = slide9.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(0.4))
p = int_title.text_frame.paragraphs[0]
p.text = "VERIFIKASI INTEGRITAS & LOGIKA"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

int_box = slide9.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2))
tf_int = int_box.text_frame
tf_int.word_wrap = True
add_bullet_points(tf_int, [
    ("Deteksi Manipulasi (SHA-256):", True, 0, ACCENT_BLUE),
    ("Pengujian manipulasi dilakukan dengan mengubah 1 karakter ciphertext di database MongoDB.", False, 1),
    ("Sistem di sisi penerima mendeteksi ketidaksesuaian nilai hash SHA-256 pasca dekripsi, menampilkan error, dan menolak merender teks pesan.", False, 1, ACCENT_RED),
    ("Pengujian Logika (White Box Testing):", True, 0, ACCENT_BLUE),
    ("Pengujian path logika klasifikasi, enkripsi, dan dekripsi pesan menghasilkan Cyclomatic Complexity yang efisien.", False, 1),
    ("Semua jalur percabangan (path) berhasil dilewati dengan tingkat keberhasilan pengujian 100%.", False, 1)
], font_size=13)


# ----------------- SLIDE 10: CONCLUSION & SUGGESTIONS -----------------
slide10 = prs.slides.add_slide(slide_layout)
set_slide_background(slide10)
add_title(slide10, "Kesimpulan & Saran Pengembangan")

# Left Column (Conclusions)
draw_card(slide10, 0.8, 1.6, 5.6, 5.1, bg_color=CARD_BG, border_color=ACCENT_GREEN)
concl_title = slide10.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(0.4))
p = concl_title.text_frame.paragraphs[0]
p.text = "KESIMPULAN PENELITIAN"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN

concl_box = slide10.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.2))
tf_con = concl_box.text_frame
tf_con.word_wrap = True
add_bullet_points(tf_con, [
    ("1. Aplikasi ChatKu berhasil mengintegrasikan Client-Side Naive Bayes dan Hybrid Cryptography AES-RSA secara E2EE murni.", False, 0),
    ("2. Klasifikasi tingkat risiko berjalan lokal di RAM HP pengirim sehingga kerahasiaan teks dan kunci terjamin penuh dari server.", False, 0),
    ("3. Model Naive Bayes memperoleh akurasi 92,90% pada data uji dengan model Multinomial NB yang efisien.", False, 0),
    ("4. Pengujian integritas SHA-256 terbukti 100% mendeteksi manipulasi database, dan audit jaringan membuktikan zero-knowledge server berfungsi penuh.", False, 0)
], font_size=13)

# Right Column (Suggestions)
draw_card(slide10, 6.9, 1.6, 5.6, 5.1, bg_color=CARD_BG, border_color=ACCENT_BLUE)
sug_title = slide10.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(0.4))
p = sug_title.text_frame.paragraphs[0]
p.text = "SARAN PENGEMBANGAN"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

sug_box = slide10.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.2), Inches(4.2))
tf_sug = sug_box.text_frame
tf_sug.word_wrap = True
add_bullet_points(tf_sug, [
    ("1. Penerapan Perfect Forward Secrecy (PFS) berbasis Diffie-Hellman agar kebocoran private key tidak membocorkan riwayat chat masa lalu.", False, 0),
    ("2. Menyeimbangkan dataset dengan menambahkan sampel kata berisiko halus atau sarkasme.", False, 0),
    ("3. Melakukan pengukuran waktu komputasi enkripsi RSA-2048 di berbagai spesifikasi hardware perangkat mobile.", False, 0),
    ("4. Menambahkan fitur QR Code / Safety Number untuk verifikasi kunci publik kontak.", False, 0),
    ("5. Mengembangkan protokol E2EE untuk grup chat menggunakan protokol Sender Keys.", False, 0)
], font_size=13)

# Save presentation
prs.save("Skripsi_Dimas_Presentation_v3.pptx")
print("Presentation created successfully as Skripsi_Dimas_Presentation_v3.pptx")
